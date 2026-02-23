from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Crea presentazione con dimensioni template
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.62)

# Colori dal template Avrai
ORANGE = RGBColor(0xFF, 0x80, 0x55)
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(0x1a, 0x1a, 0x1a)
CARD_BG = RGBColor(0x25, 0x25, 0x25)
RED_SOFT = RGBColor(0x99, 0x44, 0x44)
GREEN_SOFT = RGBColor(0x44, 0x99, 0x55)
PURPLE = RGBColor(0x90, 0x55, 0xFF)
BLUE = RGBColor(0x55, 0xB4, 0xFF)

slide_number = 0

def add_dark_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_section_slide(title_line1, title_line2="", subtitle=""):
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    # Barra laterale arancione
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    # Titolo riga 1
    y_pos = Inches(1.8) if title_line2 else Inches(2.2)
    txBox = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(8), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_line1
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    if title_line2:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(8), Inches(0.9))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_line2
        p2.font.name = "Arial"
        p2.font.size = Pt(48)
        p2.font.bold = True
        p2.font.color.rgb = ORANGE

    if subtitle:
        sub_y = Inches(3.6) if title_line2 else Inches(3.2)
        txBox3 = slide.shapes.add_textbox(Inches(1.5), sub_y, Inches(7), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = subtitle
        p3.font.name = "Arial"
        p3.font.size = Pt(16)
        p3.font.color.rgb = WHITE

    return slide

def add_header(slide, section_name):
    global slide_number
    # Numero pagina
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(0.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Nome sezione
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(3), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_name.upper()
    p2.font.name = "Arial"
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE

def add_content_slide(section_name, title, bullets):
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo slide
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = title.upper()
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    y = 1.3
    for bullet in bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {bullet}" if not bullet.startswith("•") and bullet.strip() else bullet
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        y += 0.55

    return slide

def add_4d_overview_slide(section_name):
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "LE 4D DELLA COLLABORAZIONE CON L'AI"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    d_items = [
        ("1", "Delegation", "Scegliere cosa delegare"),
        ("2", "Description", "Descrivere bene il compito"),
        ("3", "Discernment", "Valutare l'output"),
        ("4", "Diligence", "Verificare accuratezza")
    ]

    x_positions = [0.3, 2.65, 5.0, 7.35]

    for i, (num, name, desc) in enumerate(d_items):
        x = x_positions[i]
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(2.2), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(0.08), Inches(3.8))
        border.fill.solid()
        border.fill.fore_color.rgb = ORANGE
        border.line.fill.background()

        # Numero
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.6), Inches(1), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.color.rgb = ORANGE

        # Nome
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.0), Inches(1.9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Descrizione
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(1.9), Inches(2.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    return slide

def add_4d_detail_slide(section_name, d_number, d_name, d_desc, examples_good, examples_bad, question):
    """Slide dettaglio per singola D"""
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo con D
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{d_number}. {d_name.upper()}"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # Descrizione
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = d_desc
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Box SI
    card_good = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.8), Inches(4.5), Inches(2.2))
    card_good.fill.solid()
    card_good.fill.fore_color.rgb = CARD_BG
    card_good.line.fill.background()

    border_good = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.8), Inches(0.08), Inches(2.2))
    border_good.fill.solid()
    border_good.fill.fore_color.rgb = GREEN_SOFT
    border_good.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SI"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GREEN_SOFT

    y = 2.2
    for ex in examples_good:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4.2), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {ex}"
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        y += 0.4

    # Box NO
    card_bad = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.8), Inches(4.5), Inches(2.2))
    card_bad.fill.solid()
    card_bad.fill.fore_color.rgb = CARD_BG
    card_bad.line.fill.background()

    border_bad = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.8), Inches(0.08), Inches(2.2))
    border_bad.fill.solid()
    border_bad.fill.fore_color.rgb = RED_SOFT
    border_bad.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "NO"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED_SOFT

    y = 2.2
    for ex in examples_bad:
        txBox = slide.shapes.add_textbox(Inches(5.2), Inches(y), Inches(4.2), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {ex}"
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        y += 0.4

    # Domanda chiave
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'Domanda chiave: "{question}"'
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = ORANGE

    return slide

def add_three_level_comparison(section_name, title,
                               bad_prompt, bad_output,
                               good_prompt, good_output,
                               advanced_steps):
    """Slide con 3 livelli: scarso, efficace, avanzato"""
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 3 colonne
    col_width = 3.0
    col_positions = [0.3, 3.45, 6.6]
    colors = [RED_SOFT, GREEN_SOFT, PURPLE]
    labels = ["SCARSO", "EFFICACE", "AVANZATO"]
    prompts = [bad_prompt, good_prompt, "Multi-step orchestrato"]
    outputs = [bad_output, good_output, "\n".join(advanced_steps)]

    for i, (x, color, label, prompt, output) in enumerate(zip(col_positions, colors, labels, prompts, outputs)):
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.0), Inches(col_width), Inches(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        # Border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.0), Inches(0.06), Inches(4.3))
        border.fill.solid()
        border.fill.fore_color.rgb = color
        border.line.fill.background()

        # Label
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.1), Inches(col_width - 0.2), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = "Arial"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color

        # Prompt
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.4), Inches(col_width - 0.25), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{prompt}"' if i < 2 else prompt
        p.font.name = "Arial"
        p.font.size = Pt(8)
        p.font.italic = True
        p.font.color.rgb = WHITE

        # Output
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.4), Inches(col_width - 0.25), Inches(2.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = output
        p.font.name = "Arial"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(180, 180, 180)

    return slide

def add_tools_slide(section_name, title, tools):
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = title.upper()
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    y = 1.3
    for tool_name, tool_desc in tools:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = tool_name
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        txBox = slide.shapes.add_textbox(Inches(3.2), Inches(y), Inches(6.5), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tool_desc
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        y += 0.7

    return slide

def add_agenda_slide(section_name):
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "AGENDA"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    agenda = [
        ("1", "AI Fluency: Le 4D", "25 min"),
        ("2", "Esempi Testi (3 livelli)", "25 min"),
        ("3", "Esempi Immagini", "25 min"),
        ("4", "Strumenti Overview", "20 min"),
        ("5", "Esercitazione Interattiva", "20 min"),
        ("6", "Q&A", "5 min"),
    ]

    y = 1.25
    for num, item, duration in agenda:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.52))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.08), Inches(0.52))
        border.fill.solid()
        border.fill.fore_color.rgb = ORANGE
        border.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.08), Inches(0.4), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = ORANGE

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.08), Inches(6), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        txBox = slide.shapes.add_textbox(Inches(8), Inches(y + 0.08), Inches(1.3), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = duration
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.RIGHT

        y += 0.62

    return slide

def add_image_example_slide(section_name, title, description, steps):
    """Slide per esempi immagini con approccio iterativo"""
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Descrizione
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(180, 180, 180)

    # Steps come cards orizzontali
    y = 1.9
    for i, (step_title, step_desc) in enumerate(steps):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.75))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.08), Inches(0.75))
        border.fill.solid()
        border.fill.fore_color.rgb = BLUE
        border.line.fill.background()

        # Step number
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(0.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE

        # Step title
        txBox = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.1), Inches(2.5), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Step desc
        txBox = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.38), Inches(8.2), Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(180, 180, 180)

        y += 0.85

    return slide

def add_interactive_slide(section_name):
    """Slide che spiega l'esercitazione interattiva"""
    global slide_number
    slide_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, section_name)

    # Titolo
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ESERCITAZIONE INTERATTIVA"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Box con istruzioni
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()

    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(0.08), Inches(3.5))
    border.fill.solid()
    border.fill.fore_color.rgb = ORANGE
    border.line.fill.background()

    instructions = [
        "1. Scansiona il QR code con il tuo smartphone",
        "2. Scegli un prompt template dalla pagina",
        "3. Clicca 'Copia' e apri ChatGPT",
        "4. Incolla e sostituisci i campi [IN ARANCIONE]",
        "5. Premi invio e guarda il risultato!"
    ]

    y = 1.5
    for instr in instructions:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = instr
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        y += 0.55

    # Placeholder per QR
    qr_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(2), Inches(2))
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = WHITE
    qr_box.line.color.rgb = ORANGE

    txBox = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "[QR CODE]"
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Nota in basso
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Pagina con tutti i prompt: [URL DA INSERIRE]"
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = ORANGE

    return slide


# ==========================================
# CREAZIONE SLIDE
# ==========================================

# 1. Titolo
add_section_slide("AI per il", "Marketing", "Creare Contenuti e Immagini nel 2025")

# 2. Agenda
add_agenda_slide("AI Marketing")

# 3. Sezione AI Fluency
add_section_slide("AI", "Fluency", "Collaborare efficacemente con l'intelligenza artificiale")

# 4. Cos'è l'AI Fluency
add_content_slide("AI Fluency", "Cos'è l'AI Fluency?", [
    "La capacità di collaborare efficacemente con l'intelligenza artificiale",
    "",
    "Non si tratta di sapere come funziona l'AI...",
    "...ma di sapere come comunicare con l'AI",
    "",
    "L'AI non legge nel pensiero: più contesto dai, migliori risultati ottieni"
])

# 5. Le 4D Overview
add_4d_overview_slide("AI Fluency")

# 6-9. Le 4D in dettaglio
add_4d_detail_slide(
    "AI Fluency", "1", "Delegation",
    "Scegliere quali task delegare all'AI - non tutto è adatto",
    ["Bozze e prime stesure", "Brainstorming e varianti", "Ricerca e sintesi", "Task ripetitivi"],
    ["Decisioni strategiche finali", "Valutazioni etiche", "Relazioni personali", "Dati sensibili senza policy"],
    "Questo task è adatto all'AI?"
)

add_4d_detail_slide(
    "AI Fluency", "2", "Description",
    "Fornire contesto completo - l'AI non conosce la tua azienda",
    ["Target specifico con età e pain point", "Tono di voce definito", "Vincoli chiari (lunghezza, formato)", "Esempi di cosa evitare"],
    ["Prompt vaghi senza contesto", "Assumere che 'capisce'", "Non specificare il formato output", "Saltare i vincoli"],
    "Ho dato abbastanza contesto?"
)

add_4d_detail_slide(
    "AI Fluency", "3", "Discernment",
    "Valutare criticamente ogni output - non tutto è pubblicabile",
    ["Verificare tono e accuratezza", "Chiedere varianti se non convince", "Iterare e raffinare", "Usare come bozza, non finale"],
    ["Pubblicare senza leggere", "Accettare il primo output", "Ignorare incongruenze", "Fidarsi di numeri/date"],
    "Questo risultato è usabile così com'è?"
)

add_4d_detail_slide(
    "AI Fluency", "4", "Diligence",
    "Verificare fatti, numeri e fonti - l'AI può 'allucinare'",
    ["Verificare statistiche citate", "Controllare nomi e date", "Fact-check affermazioni", "Validare con fonti esterne"],
    ["Fidarsi di numeri inventati", "Citare senza verificare", "Assumere che i fatti siano corretti", "Pubblicare dati sensibili"],
    "Ho verificato i fatti?"
)

# 10. Sezione Esempi Testi
add_section_slide("Esempi", "Testi", "Scarso vs Efficace vs Avanzato")

# 11. Esempio 1: Post Instagram (3 livelli)
add_three_level_comparison(
    "Esempi Testi", "Esempio 1: Post Instagram",
    "Scrivi un post Instagram per il nostro prodotto",
    "Output:\nScopri il nostro fantastico prodotto!\nNon perdere questa occasione!\n#prodotto #novità #shopping\n\n→ Generico, zero valore",
    "Post per CloudCRM, gestionale PMI. Target: imprenditori 35-50 frustrati da Excel. Max 120 char, 3 hashtag.",
    "Output:\nGestisci i clienti su Excel? CloudCRM: tutto in un posto.\nProva gratis → link in bio\n#CRM #PMI\n\n→ Specifico, usabile",
    [
        "STEP 1: Genera 5 varianti",
        "STEP 2: Valuta (hook, CTA, engagement 1-10)",
        "STEP 3: Ottimizza la migliore",
        "",
        "→ Output validato e ottimizzato"
    ]
)

# 12. Esempio 2: Email Marketing (3 livelli)
add_three_level_comparison(
    "Esempi Testi", "Esempio 2: Email Marketing",
    "Fammi una email di marketing",
    "Output:\nOggetto: Offerta speciale!\nCaro cliente, non perdere...\n\n→ Spam folder garantito",
    "Email riattivazione B2B, clienti inattivi 6 mesi. Software gestionale, sconto 20%. Max 100 parole, no urgenza finta.",
    "Output:\nOggetto: Il tuo gestionale ti aspetta\nSono 6 mesi. Abbiamo aggiunto...\nRinnova con -20%\n\n→ Personalizzato, credibile",
    [
        "STEP 1: Crea sequenza 3 email",
        "STEP 2: 2 varianti oggetto (A/B)",
        "STEP 3: Predici open rate",
        "",
        "→ Sequenza completa + test"
    ]
)

# 13. Esempio 3: Brainstorming (3 livelli)
add_three_level_comparison(
    "Esempi Testi", "Esempio 3: Brainstorming",
    "Dammi idee per una campagna",
    "Output:\n1. Social media\n2. Influencer\n3. Contest\n4. Email\n\n→ Lista generica",
    "5 concept social B2C. Abbigliamento sportivo sostenibile. Donne 25-40. Per ogni: nome, visual, copy, canali.",
    "Output:\n1. #RiVesti - UGC riciclo\n2. 30GiorniGreen - Challenge\n...\n\n→ Concept actionable",
    [
        "STEP 1: Analizza competitor settore",
        "STEP 2: 3 concept differenzianti",
        "STEP 3: Stress test + rischi",
        "STEP 4: Raccomandazione finale",
        "→ Strategia completa"
    ]
)

# 14. Sezione Immagini
add_section_slide("Esempi", "Immagini", "Approccio iterativo e incrementale")

# 15. Intro Immagini - 5 funzionalità
add_content_slide("Immagini AI", "Le 5 Funzionalità Chiave", [
    "Cambia Colore - modificare colori specifici (lacci, suola, etc.)",
    "Cambia Dettaglio - aggiungere/modificare elementi",
    "Foto Emozionale - creare lifestyle shots",
    "Cambia Posa - stesso soggetto, posa diversa",
    "Cambia Background - nuovo ambiente, stesso soggetto"
])

# 16. Approccio Iterativo
add_image_example_slide(
    "Immagini AI", "Approccio Iterativo",
    "Non chiedere tutto insieme. Procedi per gradi, validando ogni step.",
    [
        ("Input iniziale", "Carica l'immagine di partenza (es. 4 viste prodotto)"),
        ("Modifica singola", "Cambia UN solo attributo (es. colore lacci a #48cae4)"),
        ("Validazione", "Controlla il risultato. Se OK, usa come nuovo input"),
        ("Iterazione", "Ripeti: modifica successiva sul risultato validato"),
        ("Output finale", "Dopo 3-4 step: prodotto completamente personalizzato")
    ]
)

# 17. Esempio Cambio Colore
add_image_example_slide(
    "Immagini AI", "Esempio: Cambio Colore",
    "Modifica precisa di un singolo elemento mantenendo tutto il resto identico.",
    [
        ("Prompt", "\"Cambia il colore dei lacci da bianco a #48cae4\""),
        ("Vincoli", "Mantieni ESATTAMENTE inquadratura, luce, sfondo"),
        ("Verifica", "Il risultato deve sembrare una foto reale"),
        ("Iterazione", "Se OK → usa come base per la modifica successiva")
    ]
)

# 18. Esempio Catena Completa
add_image_example_slide(
    "Immagini AI", "Esempio: Catena di 3 Modifiche",
    "Trasforma un campione in variante complessa con controllo totale.",
    [
        ("Step 1", "Lacci: bianco → blu (#48cae4) → Valida"),
        ("Step 2", "Suola: beige → verde (#a7c957) → Valida"),
        ("Step 3", "Dettaglio tallone: aggiungi pelle arancione (#ff9e00)"),
        ("Risultato", "3 dettagli modificati, massimo realismo fotografico")
    ]
)

# 19. Esempio Foto Emozionale
add_image_example_slide(
    "Immagini AI", "Esempio: Foto Lifestyle",
    "Crea immagini emozionali per campagne advertising.",
    [
        ("Soggetto", "Donna 30 anni, espressione energica, in movimento"),
        ("Prodotto", "Indossa/usa il prodotto in modo naturale"),
        ("Location", "Parco urbano, luce naturale mattutina"),
        ("Variante", "Cambia solo background: stesso soggetto, location diversa")
    ]
)

# 20. Sezione Strumenti
add_section_slide("Gli", "Strumenti", "Panoramica dei tool disponibili")

# 21. Strumenti Testi
add_tools_slide("Strumenti", "Creazione Testi", [
    ("ChatGPT", "Il più versatile - copy, email, brainstorming"),
    ("Claude", "Testi lunghi e ragionamento complesso"),
    ("Gemini", "Integrato con Google Workspace"),
    ("Jasper", "Specifico marketing - template pronti"),
    ("Copy.ai", "Workflow collaborativi per team"),
])

# 22. Strumenti Immagini
add_tools_slide("Strumenti", "Creazione Immagini", [
    ("Midjourney", "Qualità artistica top - via Discord"),
    ("DALL-E 3", "Dentro ChatGPT - facile e veloce"),
    ("Adobe Firefly", "Integrato con Photoshop - diritti chiari"),
    ("Canva AI", "Per chi già usa Canva"),
    ("Ideogram", "Ottimo per testo nelle immagini"),
])

# 23. Strumenti Video
add_tools_slide("Strumenti", "Video e Audio", [
    ("Runway", "Video da testo/immagini"),
    ("HeyGen", "Avatar parlanti"),
    ("ElevenLabs", "Voci sintetiche realistiche"),
    ("Suno", "Musica e jingle da testo"),
    ("CapCut", "Editing video con AI"),
])

# 24. Sezione Esercitazione
add_section_slide("Ora", "Provate Voi!", "Esercitazione interattiva")

# 25. Slide interattiva
add_interactive_slide("Esercitazione")

# 26. Tips
add_content_slide("Tips", "Per Iniziare Domani", [
    "",
    "Parti da UN solo strumento e imparalo bene",
    "",
    "Salva i prompt che funzionano in un documento condiviso",
    "",
    "Non fidarti ciecamente: verifica sempre fatti e numeri",
    "",
    "Itera: il primo output raramente è quello finale"
])

# 27. Chiusura
add_section_slide("Domande?", "", "Grazie per l'attenzione!")


# Salva
output_path = "/Users/gianpaologreco/Masterclass_AI_Marketing/Masterclass_AI_Marketing_v3.pptx"
prs.save(output_path)
print(f"Presentazione salvata: {output_path}")
print(f"Totale slide: {slide_number}")
